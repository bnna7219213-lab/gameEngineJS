#!/usr/bin/env python3
"""生成《宣传+操作手册.pptx》：engine_tensorflow+js 编辑器 Debug 模式（边调试边热加载边运行）。
依赖 python-pptx（已确认可用）。纯文档产物，不影响引擎逻辑。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ACCENT = RGBColor(0x44, 0x72, 0xC4)
ACCENT2 = RGBColor(0xED, 0x7D, 0x31)
DARK = RGBColor(0x22, 0x26, 0x2E)
GREY = RGBColor(0x5A, 0x60, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG = RGBColor(0xF4, 0xF6, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_band(slide, color=ACCENT, h=Inches(1.15)):
    box = slide.shapes.add_shape(1, 0, 0, SW, h)  # rectangle
    box.fill.solid(); box.fill.fore_color.rgb = color
    box.line.fill.background()
    box.shadow.inherit = False
    return box


def add_text(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = ln.get('align', align)
        if 'space_before' in ln: p.space_before = Pt(ln['space_before'])
        if 'space_after' in ln: p.space_after = Pt(ln['space_after'])
        for run in ln['runs']:
            r = p.add_run(); r.text = run['text']
            r.font.size = Pt(run.get('size', 18))
            r.font.bold = run.get('bold', False)
            r.font.italic = run.get('italic', False)
            r.font.name = run.get('name', 'Microsoft YaHei')
            r.font.color.rgb = run.get('color', DARK)
    return tb


def slide_title(title, subtitle=None, color=ACCENT):
    s = prs.slides.add_slide(BLANK)
    add_band(s, color)
    add_text(s, Inches(0.6), Inches(0.18), SW - Inches(1.2), Inches(0.8),
             [{'runs': [{'text': title, 'size': 30, 'bold': True, 'color': WHITE}]}],
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(s, Inches(0.6), Inches(1.2), SW - Inches(1.2), Inches(0.5),
                 [{'runs': [{'text': subtitle, 'size': 14, 'color': GREY}]}])
    # background tint
    bg = s.shapes.add_shape(1, 0, Inches(1.15), SW, SH - Inches(1.15))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)  # send to back
    return s


def bullets(slide, items, top=Inches(1.5), left=Inches(0.7), size=18, gap=10, color=DARK):
    lines = []
    for it in items:
        if isinstance(it, tuple):
            head, sub = it
            lines.append({'runs': [{'text': '• ' + head, 'size': size, 'bold': True, 'color': ACCENT}], 'space_after': 2, 'space_before': gap})
            lines.append({'runs': [{'text': '    ' + sub, 'size': size - 3, 'color': color}], 'space_after': gap})
        else:
            lines.append({'runs': [{'text': '• ' + it, 'size': size, 'color': color}], 'space_after': gap, 'space_before': gap})
    # place above band: ensure after background; add a white textbox on top
    tb = add_text(slide, left, top, SW - left - Inches(0.5), SH - top - Inches(0.4), lines)
    return tb


# ---------- 1. 封面 ----------
s = prs.slides.add_slide(BLANK)
band = add_band(s, DARK, h=SH)
band.fill.fore_color.rgb = DARK
add_text(s, Inches(0.9), Inches(2.2), SW - Inches(1.8), Inches(1.4),
         [{'runs': [{'text': 'engine_tensorflow+js', 'size': 46, 'bold': True, 'color': WHITE}]}],
         anchor=MSO_ANCHOR.MIDDLE)
add_text(s, Inches(0.9), Inches(3.5), SW - Inches(1.8), Inches(1.2),
         [{'runs': [{'text': '自研引擎 · 浏览器端 · 计算与推理渲染', 'size': 22, 'color': RGBColor(0x9D,0xC3,0xFF)}]},
          {'runs': [{'text': '代码工作台 Debug 模式：边调试 · 边热加载 · 边运行', 'size': 18, 'color': RGBColor(0xED,0x7D,0x31), 'bold': True}], 'space_before': 10}])
add_text(s, Inches(0.9), Inches(6.4), SW - Inches(1.8), Inches(0.5),
         [{'runs': [{'text': '宣传 + 操作手册', 'size': 14, 'color': RGBColor(0xAA,0xB0,0xBC)}]}])

# ---------- 2. 产品定位 ----------
s = slide_title('一句话定位', 'C++17 引擎的全功能对等 JS 重构版', color=ACCENT)
bullets(s, [
    ('渲染', '自研引擎：WebGPU(L3) → WebGL2(L1/L2) → Software(L0 黄金参考) 三级后端，无需 three.js。'),
    ('计算推理', 'TensorFlow.js 浏览器端推理（神经材质 / DDGI 去噪 / 超分 / AI 策略）；缺失时降级内置 NanoTensor。'),
    ('编辑器', 'Web 版 IDE（原生 DOM，无框架、无构建步骤），功能对等 python/ide。'),
    ('游戏', 'games/ 下 Breakout / Space Shooter / Match3 / Action3D 浏览器直接运行。'),
], top=Inches(1.5))

# ---------- 3. 核心能力总览 ----------
s = slide_title('编辑器核心能力', '原生 DOM 实现，红线 D/E/F 全部遵守', color=ACCENT2)
bullets(s, [
    '3D 视窗：相机轨道/平移、多选、Gizmo 平移-旋转-缩放、吸附步进、撤销/重做、聚焦、AABB 拾取。',
    '层级 / 检视器：组件注册表驱动，PBR 参数、灯光、ECS、碰撞体等可编辑。',
    '时间轴：关键帧动画、轨道、播放/停止/拖拽。',
    '代码工作台：自研语法高亮 + 文件树 + OPFS 持久化 + Web Worker 语法校验。',
    'Profiler：分段计时 / drawCalls / 内存 / 帧时间线可视化（60/30fps 参考线）。',
    'Play 模式：场景深拷贝快照运行，停止即回滚，编辑态零污染（红线 F）。',
], top=Inches(1.5), size=16, gap=8)

# ---------- 4. Debug 模式（宣传重点） ----------
s = slide_title('Debug 模式：边调试 · 边热加载 · 边运行', '你改代码，游戏不重启，立即看到效果', color=ACCENT)
bullets(s, [
    ('边运行', '点 ▶Run（Ctrl+P）进入 Play，游戏在独立快照上实时运行。'),
    ('边调试', '代码工作台底部工具条：⏸暂停/▶继续、⏭单步（精确 1/60s 一帧）、🔄热重载；实时面板显示每个实体 pos(x,y,z) 与脚本日志。'),
    ('边热加载', '运行中改脚本按 Ctrl+S → 新代码立即热加载进运行会话（不停止游戏），并写回编辑器对象脚本（下次 Play 也生效）。'),
    ('安全', '语法错误被拦截并提示，旧脚本继续运行不中断；单脚本异常被隔离，不影响其它脚本与编辑数据。'),
], top=Inches(1.5))

# ---------- 5. 操作手册：步骤 ----------
s = slide_title('操作手册：5 步上手 Debug 模式', '从零跑通「边调试边热加载边运行」', color=ACCENT2)
bullets(s, [
    '1) 启动：node serve.mjs → 浏览器打开 http://localhost:8080/ 进入 src/editor。',
    '2) 进入 Play：点工具栏「Play 模式」按钮（或按 Ctrl+P），状态栏显示「Play 模式」。',
    '3) 打开代码工作台：点工具栏「代码工作台」按钮，默认打开 scripts/cube.js。',
    '4) 热加载：把 this.move(0.5*dt,0,0) 改成别的逻辑，按 Ctrl+S → 游戏立即按新脚本运行。',
    '5) 调试：用 ⏸暂停 冻结画面，⏭单步 逐帧观察实体坐标变化，🔄热重载 重新载入当前脚本；脚本里用 this.log(...) 输出日志到面板。',
], top=Inches(1.5), size=15, gap=9)

# ---------- 6. 操作手册：脚本 API 速查 ----------
s = slide_title('操作手册：脚本 API 速查', '在脚本内以 this 调用；可用 dt(帧时长) 与 input(按键/鼠标)', color=ACCENT)
bullets(s, [
    ('移动/旋转', 'this.move(dx,dy,dz) · this.rotate(rx,ry,rz) · this.setPosition(x,y,z) · this.setRotation(x,y,z)'),
    ('读取', 'this.position / this.rotation / this.scale（数组）· this.time（运行时累计秒）。'),
    ('日志', 'this.log(msg) → 输出到代码工作台「脚本日志」面板，便于边运行边观察。'),
    ('文件约定', 'scripts/<对象名>[#槽位].js 对应同名对象的脚本槽位，如 scripts/cube.js；多脚本用 #0/#1。空槽位热加载会自动新增首个脚本。'),
    ('示例', 'this.move(0.5*dt, 0, 0);  // 每帧沿 x 移动；dt 让速度与时长无关'),
], top=Inches(1.5), size=15, gap=8)

# ---------- 7. 架构与安全保证 ----------
s = slide_title('架构红线：计算 100% 浏览器内', '性能与正确性可证，运行不依赖外部服务', color=ACCENT2)
bullets(s, [
    ('D. CPU 参考先行', '任何 GPU 路径先有 CPU 参考，Software 为黄金基准，永不下线。'),
    ('E. 可选层缺失即降级', 'WebGPU / WebGL2 / TF.js / Worker 全部可选；缺失时自动回退不崩溃。'),
    ('F. 编辑器→运行时单向', '运行时绝不反向改写编辑场景；Play 停止即丢弃快照，编辑态原封不动。'),
    ('快照隔离', 'Play 进入时对 Scene3D 做序列化往返深拷贝，物理只回写运行时实体表。'),
    ('Worker 化', '脚本语法校验在后台线程，不阻塞主线程交互。'),
], top=Inches(1.5))

# ---------- 8. 结语 / 快速开始 ----------
s = slide_title('快速开始 & 验收', '开箱即用，零外部依赖', color=ACCENT)
bullets(s, [
    ('运行', 'node serve.mjs （或 python -m http.server 8080）。'),
    ('验收', 'node tools/run_smoke.js —— 全量 smoke 必须全绿（含 debug_mode / script_compiler / file_system / profiler_timeline）。'),
    ('文档', 'README.md（含 Debug 模式章节）· CONTRACT.md（并行开发唯一事实来源）· docs/PORTING.md（C++→JS 映射）。'),
    ('配套脚本', 'tools/make_manual_pptx.py 可重新生成本手册。'),
], top=Inches(1.5), size=16, gap=9)

# ---------- Debug 工作流（流程图） ----------
def add_flow_slide(steps, caption, notes):
    s = prs.slides.add_slide(BLANK)
    add_band(s, ACCENT2, h=Inches(1.15))
    add_text(s, Inches(0.6), Inches(0.18), SW - Inches(1.2), Inches(0.8),
             [{'runs': [{'text': 'Debug 工作流：暂停 → 改脚本 → Ctrl+S → 游戏中变化', 'size': 26, 'bold': True, 'color': WHITE}]}],
             anchor=MSO_ANCHOR.MIDDLE)
    # 背景
    bg = s.shapes.add_shape(1, 0, Inches(1.15), SW, SH - Inches(1.15))
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); bg.shadow.inherit = False
    s.shapes._spTree.remove(bg._element); s.shapes._spTree.insert(2, bg._element)
    n = len(steps)
    bw, bh = Inches(2.7), Inches(1.5)
    gap = Inches(0.55)
    total = n * bw + (n - 1) * gap
    x0 = (SW - total) / 2
    top = Inches(2.2)
    for i, txt in enumerate(steps):
        x = x0 + i * (bw + gap)
        box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, bw, bh)
        box.fill.solid(); box.fill.fore_color.rgb = ACCENT if i % 2 == 0 else ACCENT2
        box.line.color.rgb = WHITE
        tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for j, line in enumerate(txt):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = line; r.font.size = Pt(15 if j == 0 else 13); r.font.bold = (j == 0)
            r.font.color.rgb = WHITE; r.font.name = 'Microsoft YaHei'
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + bw + Inches(0.05), top + bh / 2 - Inches(0.18), gap - Inches(0.1), Inches(0.36))
            ar.fill.solid(); ar.fill.fore_color.rgb = GREY; ar.line.fill.background(); ar.shadow.inherit = False
    add_text(s, Inches(0.8), Inches(4.2), SW - Inches(1.6), Inches(0.5),
             [{'runs': [{'text': caption, 'size': 15, 'bold': True, 'color': DARK}]}])
    bullets(s, notes, top=Inches(4.7), size=14, gap=7)


add_flow_slide(
    steps=[['⏸ 暂停', '（或边运行）'], ['✎ 改脚本', '代码工作台'], ['Ctrl+S', '热加载到运行会话'], ['▶ 游戏中', '实时生效']],
    caption='一句话流程：游戏不停，改完即见。',
    notes=[
        '1) 进入 Play（Ctrl+P）后，点 ⏸暂停 冻结画面（也可边运行边改）。',
        '2) 在代码工作台（或检视器脚本卡）编辑对象脚本。',
        '3) 按 Ctrl+S：新代码立即热加载进运行会话，不重启游戏；语法错误被拦截，旧脚本继续跑。',
        '4) 切回 ▶继续 或单步观察：实体状态、this.log 日志、监视表达式实时刷新。',
        '进阶：脚本内 this.breakpoint() 命中点自动暂停；⏭单步 逐帧观察；实时面板显示每个实体 pos 与脚本输出。',
    ],
)

out = '宣传+操作手册.pptx'
prs.save(out)
# 自检：用同一 Python 字符串重新打开，确认 OOXML 包可被 python-pptx 解析（绕过 shell 对中文路径的编码问题）
chk = Presentation(out)
print('saved', out, 'slides=', len(prs.slides._sldIdLst), 'reopened-ok=', len(chk.slides._sldIdLst))
